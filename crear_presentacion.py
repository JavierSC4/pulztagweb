from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import re

# Función para verificar si una imagen existe
def verificar_imagen(ruta_imagen):
    if not os.path.isfile(ruta_imagen):
        print(f"Advertencia: La imagen '{ruta_imagen}' no se encontró.")
        return False
    return True

# Función para procesar y agregar texto con negrita
def agregar_texto_con_negrita(text_frame, text):
    segments = parse_bold_text(text)
    for segment, bold in segments:
        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = segment
        run.font.bold = bold
        run.font.name = 'Lato'
        run.font.size = Pt(20)
        run.font.color.rgb = AZUL_NAVIO
        p.alignment = PP_ALIGN.LEFT

# Crear una nueva presentación
prs = Presentation()

# Definir colores en RGB
AZUL_NAVIO = RGBColor(44, 62, 80)        # #2C3E50
TURQUESA = RGBColor(26, 188, 156)        # #1ABC9C
BLANCO = RGBColor(255, 255, 255)         # #FFFFFF
GRIS_CLARO = RGBColor(236, 240, 241)     # #ECF0F1
DORADO_SUAVE = RGBColor(241, 196, 15)    # #F1C40F

# Ruta de las imágenes
IMAGES_DIR = "images"
LOGO_PATH = os.path.join(IMAGES_DIR, "logo.png")
PRODUCT_IMAGES = {
    "product1": os.path.join(IMAGES_DIR, "product1.jpg"),
    "product2": os.path.join(IMAGES_DIR, "product2.jpg"),
    "product3": os.path.join(IMAGES_DIR, "product3.jpg"),
    "product4": os.path.join(IMAGES_DIR, "product4.jpg"),
    "product5": os.path.join(IMAGES_DIR, "product5.jpg"),
}
BACKGROUND_IMAGE = os.path.join(IMAGES_DIR, "background1.jpg")

# Función para agregar una diapositiva con título y contenido
def agregar_diapositiva_titulo_con_contenido(prs, titulo, contenido, imagen=None):
    slide_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)
    
    # Añadir logotipo si existe
    if verificar_imagen(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(9), Inches(0.2), width=Inches(1.5))
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = ""  # Limpiar texto inicial
    run = p.add_run()
    run.text = titulo
    run.font.name = 'Montserrat'
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = AZUL_NAVIO
    p.alignment = PP_ALIGN.LEFT
    
    # Contenido
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    agregar_texto_con_negrita(tf, contenido)
    
    # Insertar imagen si existe
    if imagen and verificar_imagen(imagen):
        try:
            left = Inches(6.8)
            top = Inches(2)
            height = Inches(3.5)
            slide.shapes.add_picture(imagen, left, top, height=height)
        except Exception as e:
            print(f"Error al insertar la imagen {imagen}: {e}")

# Función para agregar una diapositiva con título y viñetas
def agregar_diapositiva_titulo_con_viñetas(prs, titulo, viñetas, imagen=None):
    slide_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)
    
    # Añadir logotipo si existe
    if verificar_imagen(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(9), Inches(0.2), width=Inches(1.5))
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = ""  # Limpiar texto inicial
    run = p.add_run()
    run.text = titulo
    run.font.name = 'Montserrat'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = AZUL_NAVIO
    p.alignment = PP_ALIGN.LEFT
    
    # Viñetas
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for viñeta in viñetas:
        agregar_texto_con_negrita(tf, viñeta)
    
    # Insertar imagen si existe
    if imagen and verificar_imagen(imagen):
        try:
            left = Inches(7)
            top = Inches(1.5)
            height = Inches(3.5)
            slide.shapes.add_picture(imagen, left, top, height=height)
        except Exception as e:
            print(f"Error al insertar la imagen {imagen}: {e}")

# Función para agregar una diapositiva con título y lista numerada
def agregar_diapositiva_titulo_con_lista(prs, titulo, lista):
    slide_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)
    
    # Añadir logotipo si existe
    if verificar_imagen(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(9), Inches(0.2), width=Inches(1.5))
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = ""  # Limpiar texto inicial
    run = p.add_run()
    run.text = titulo
    run.font.name = 'Montserrat'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = AZUL_NAVIO
    p.alignment = PP_ALIGN.LEFT
    
    # Lista Numerada
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(lista, 1):
        # Añadir numeración manualmente
        texto = f"{idx}. {item}"
        agregar_texto_con_negrita(tf, texto)

# Función para agregar una diapositiva con título, viñetas y una imagen
def agregar_diapositiva_con_imagen(prs, titulo, viñetas, imagen):
    slide_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)
    
    # Añadir logotipo si existe
    if verificar_imagen(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(9), Inches(0.2), width=Inches(1.5))
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(6), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = ""  # Limpiar texto inicial
    run = p.add_run()
    run.text = titulo
    run.font.name = 'Montserrat'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = AZUL_NAVIO
    p.alignment = PP_ALIGN.LEFT
    
    # Viñetas
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for viñeta in viñetas:
        agregar_texto_con_negrita(tf, viñeta)
    
    # Imagen
    if verificar_imagen(imagen):
        try:
            left = Inches(7)
            top = Inches(1.5)
            height = Inches(3.5)
            slide.shapes.add_picture(imagen, left, top, height=height)
        except Exception as e:
            print(f"Error al insertar la imagen {imagen}: {e}")

# Función para procesar texto con negrita
def parse_bold_text(text):
    """
    Procesa el texto y retorna una lista de tuplas (segmento, negrita),
    donde 'segmento' es una parte del texto y 'negrita' es un booleano
    indicando si el segmento debe estar en negrita.
    """
    # Patrón para encontrar texto entre **
    pattern = re.compile(r'\*\*(.*?)\*\*')
    segments = []
    last_end = 0

    for match in pattern.finditer(text):
        start, end = match.span()
        # Texto antes de **
        if start > last_end:
            segments.append((text[last_end:start], False))
        # Texto entre **
        segments.append((match.group(1), True))
        last_end = end

    # Texto después de la última **
    if last_end < len(text):
        segments.append((text[last_end:], False))
    
    return segments

# Función para agregar texto con negrita
def agregar_texto_con_negrita(text_frame, text):
    segments = parse_bold_text(text)
    for segment, bold in segments:
        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = segment
        run.font.bold = bold
        run.font.name = 'Lato'
        run.font.size = Pt(20)
        run.font.color.rgb = AZUL_NAVIO
        p.alignment = PP_ALIGN.LEFT

# Diapositiva 1 - INTRODUCCIÓN
titulo1 = "PulztagWeb: Mejorando la Experiencia del Cliente con Tecnología Interactiva AIDC"
contenido1 = (
    "PulztagWeb mejora la experiencia del cliente mediante el uso de tecnologías interactivas AIDC "
    "(Identificación y Captura Automática de Datos) como **NFC, códigos QR, RFID y biometría**. "
    "Nuestra **plataforma web robusta** integra estas tecnologías para **optimizar** y **personalizar** cada interacción, "
    "ofreciendo **soluciones innovadoras** que conectan eficazmente a las empresas con sus clientes."
)
agregar_diapositiva_titulo_con_contenido(prs, titulo1, contenido1, imagen=None)

# Diapositiva 2 – PROPUESTA DE VALOR
titulo2 = "Nuestra Propuesta de Valor"
viñetas2 = [
    "**Tecnología Avanzada:** Utilizamos tecnologías AIDC integradas para extraer el máximo valor de cada interacción con tus clientes.",
    "**Personalización Real:** Nuestro modelo único de **CEM** permite una configuración intuitiva de puntos de acceso (con AIDC) para ser utilizados como:\n"
    "   - Encuestas de satisfacción\n"
    "   - Puntos informativos\n"
    "   - Catálogos de productos\n"
    "   - Función \"Añade a tu carrito\"\n"
    "   - Descripción de inventario, entre otros.",
    "**Resultados Medibles:** Proporcionamos métricas claras para apoyar decisiones estratégicas, incluyendo:\n"
    "   - **Net Promoter Score (NPS)**\n"
    "   - **Customer Satisfaction Score (CSAT)**\n"
    "   - **Customer Effort Score (CES)**\n"
    "   - **Tasa de Retención de Clientes**\n"
    "   - **Churn Rate (Tasa de Deserción)**\n"
    "   - **Tiempo de Respuesta y Resolución**\n"
    "   - **Customer Lifetime Value**\n"
    "   - **First Contact Resolution**\n"
    "   - **Sentiment Score**\n"
    "   - **Tasa de Resolución de Quejas**\n"
    "   - **Customer Engagement Score**\n"
    "   - **Revenue Per Customer**"
]
agregar_diapositiva_titulo_con_viñetas(prs, titulo2, viñetas2, imagen=PRODUCT_IMAGES.get("product1"))

# Diapositiva 3 – Nuestra Misión y Visión
titulo3 = "Nuestra Misión y Visión"
viñetas3 = [
    "**Misión:** Transformar la experiencia de interacción entre empresas y clientes mediante soluciones AIDC personalizadas que **optimicen procesos**, **generen conexiones significativas** y **mejoren la eficiencia operativa**.",
    "**Visión:** Ser la **empresa líder** en **soluciones AIDC innovadoras** en **Latinoamérica**, facilitando la **transformación digital** y creando un mundo más **interactivo, eficiente y conectado**."
]
agregar_diapositiva_titulo_con_lista(prs, titulo3, viñetas3)

# Diapositiva 4 – Desafíos Actuales
titulo4 = "Desafíos Actuales en la Experiencia del Cliente"
viñetas4 = [
    "**Interacciones Fragmentadas:** Dificultad para gestionar múltiples canales de comunicación.",
    "**Falta de Personalización:** Experiencias genéricas que no satisfacen las necesidades específicas de los clientes, lo cual es común en el mercado.",
    "**Medición de Resultados:** Ausencia de métricas claras para evaluar la efectividad de las estrategias.",
    "**Eficiencia Operativa:** Procesos manuales que consumen tiempo y recursos, como gráficos físicos, códigos QR y encuestas de satisfacción."
]
agregar_diapositiva_titulo_con_viñetas(prs, titulo4, viñetas4, imagen=PRODUCT_IMAGES.get("product2"))

# Diapositiva 5 - Nuestra Solución: Plataforma CEM de PulztagWeb
titulo5 = "Nuestra Solución: Plataforma CEM de PulztagWeb"
contenido5 = (
    "Una **Plataforma Integral de Gestión de la Experiencia del Cliente (CEM)** que centraliza todas las herramientas necesarias para **gestionar, optimizar y personalizar** cada interacción con tus clientes en una única plataforma."
)
agregar_diapositiva_titulo_con_contenido(prs, titulo5, contenido5, imagen=None)

# Diapositiva 6 - Características Principales
titulo6 = "Características Principales de Nuestra Plataforma"
viñetas6 = [
    "**Administración de URLs:** Gestiona URLs dedicadas para atención al cliente, acceso a productos y evaluación de servicios.",
    "**Gestión de Pulzcards:** Crea y administra tarjetas virtuales con información de contacto y códigos QR personalizados.",
    "**Gestión de Etiquetas (Tags):** Administra etiquetas NFC y códigos QR con funcionalidades avanzadas.",
    "**Gestión de Inventario y Productos:** Organiza y gestiona tus productos y recursos de manera intuitiva.",
    "**Análisis de Datos:** Herramientas integradas para analizar el comportamiento y preferencias de tus clientes.",
    "**Fidelización de Clientes:** Implementa programas de lealtad efectivos para aumentar la fidelidad.",
    "**Integración de Sistemas:** Integramos nuestras soluciones con tus sistemas existentes para optimizar procesos.",
    "**Soporte y Mantenimiento:** Ofrecemos soporte técnico continuo para asegurar el funcionamiento óptimo."
]
agregar_diapositiva_titulo_con_viñetas(prs, titulo6, viñetas6, imagen=PRODUCT_IMAGES.get("product3"))

# Diapositiva 7 - Beneficios Clave para tu Negocio
titulo7 = "Beneficios Clave para tu Negocio"
viñetas7 = [
    "**Tecnología Avanzada:** Integración fluida de **NFC, QR, RFID y biometría** con respaldo de **IA y AIDC**.",
    "**Personalización Real:** Adaptación a cada cliente, sector y necesidad específica.",
    "**Resultados Medibles:** Optimización de procesos y obtención de métricas claras para decisiones estratégicas.",
    "**Centralización de Información:** Consolidación de datos e interacciones en una única plataforma.",
    "**Eficiencia Operativa:** Automatización de procesos que reduce costos y tiempo operativo.",
    "**Seguridad de Datos:** Protección avanzada de la información de tus clientes, garantizando privacidad e integridad."
]
agregar_diapositiva_titulo_con_viñetas(prs, titulo7, viñetas7, imagen=PRODUCT_IMAGES.get("product4"))

# Diapositiva 8 - Nuestros Productos Destacados
titulo8 = "Nuestros Productos Destacados"
viñetas8 = [
    "1. **Etiqueta de Sticker de RR.SS. NFC Antimetal:**\n   - Stickers personalizables con tecnología NFC, resistentes al agua y duraderos.",
    "2. **Tarjetas PVC Inteligentes NFC:**\n   - Tarjetas inteligentes con NFC para compartir URLs, redes sociales y más.",
    "3. **Pulseras Tejidas NFC Personalizables:**\n   - Pulseras ecológicas con NFC ideales para eventos y promociones.",
    "4. **Soporte de Menú de Acrílico con QR y NFC:**\n   - Soportes resistentes al agua para menús interactivos en restaurantes y cafeterías.",
    "5. **Tarjetas de Madera NFC Empresariales:**\n   - Tarjetas de presentación ecológicas con NFC para conexiones digitales."
]
agregar_diapositiva_titulo_con_lista(prs, titulo8, viñetas8)

# Diapositiva 9 - Nuestro Compromiso
titulo9 = "Nuestro Compromiso"
viñetas9 = [
    "**Innovación Continua:** Nos mantenemos a la vanguardia tecnológica para ofrecer soluciones siempre actualizadas.",
    "**Soporte Personalizado:** Brindamos atención y soporte técnico adaptado a las necesidades de cada cliente.",
    "**Calidad y Seguridad:** Implementamos altos estándares de calidad y seguridad en todas nuestras soluciones."
]
agregar_diapositiva_titulo_con_viñetas(prs, titulo9, viñetas9, imagen=PRODUCT_IMAGES.get("product5"))

# Diapositiva 10 - Hacia Dónde Queremos Llegar
titulo10 = "Hacia Dónde Queremos Llegar"
viñetas10 = [
    "**Expansión Regional:** Consolidarnos como líderes en soluciones AIDC en Latinoamérica.",
    "**Desarrollo de Nuevas Tecnologías:** Integrar tecnologías emergentes para enriquecer la experiencia del cliente.",
    "**Ampliación de Servicios:** Incorporar nuevos servicios que respondan a las dinámicas cambiantes del mercado.",
    "**Alianzas Estratégicas:** Establecer colaboraciones con otras empresas tecnológicas para potenciar nuestras soluciones."
]
agregar_diapositiva_titulo_con_viñetas(prs, titulo10, viñetas10, imagen=BACKGROUND_IMAGE)

# Diapositiva 11 - Invitación a Transformar tu Negocio
titulo11 = "Invitación a Transformar tu Negocio"
contenido11 = (
    "Únete a las empresas que ya están optimizando sus interacciones con clientes gracias a nuestra **Plataforma CEM**. "
    "Descubre cómo **PulztagWeb** puede ayudarte a:\n\n"
    "• **Mejorar la Satisfacción del Cliente**\n"
    "• **Aumentar la Fidelidad y Retención**\n"
    "• **Optimizar Procesos Internos**\n"
    "• **Tomar Decisiones Basadas en Datos**"
)
agregar_diapositiva_titulo_con_contenido(prs, titulo11, contenido11, imagen=BACKGROUND_IMAGE)

# Guardar la presentación
prs.save('PulztagWeb_Presentacion_Elegante.pptx')

print("Presentación 'PulztagWeb_Presentacion_Elegante.pptx' creada exitosamente.")